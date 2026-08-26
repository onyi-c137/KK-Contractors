from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Slide Data
slides_data = [
    {
        "title": "Environmental Analysis & Market Entry Strategy: Even Realities G2 AR Glasses",
        "subtitle": "Task 1-A | Marketing in a Digital Era (MKT2002S) — UCD BBS45 Singapore\nPresenters: Group Presentation (5 Members)",
        "bullets": []
    },
    {
        "title": "A. Company Overview",
        "subtitle": "",
        "bullets": [
            "Founded: 2023, Shenzhen | HQ: Shenzhen + Berlin | 400+ Employees | US$150M Raised",
            "Founder: Will Wang (ex-Apple, ex-JMGO CPO)",
            "Investors: Backed by Tencent and Meituan",
            "Design Philosophy: Eyewear-first design (looks like ordinary glasses, not bulky tech)",
            "Product Evolution (G1 → G2):",
            "  • Display: +75% larger display, 50% sharper HUD",
            "  • Controls: R1 smart ring gesture control added",
            "  • Build: 36g ultralight titanium/magnesium build, IP65 rating",
            "Sources: Dealroom.co (2026); The Spectacle Factory (2025)"
        ]
    },
    {
        "title": "A. Features & Market Position",
        "subtitle": "",
        "bullets": [
            "Display: 3D floating micro-LED HUD",
            "Privacy Architecture: No camera, no speakers (strict client confidentiality design)",
            "App Suite: Real-time translation, Teleprompt, Even AI integration",
            "Key Benefits: Hands-free productivity, real-time live translation, all-day comfort",
            "Market Context: Global smart glasses market growing from US$7.53B (2026) to US$11.37B (2030)",
            "Positioning: High-trust corporate productivity tool for legal, financial, and PMET professionals",
            "Sources: Grand View Research (2026); CES Innovation Awards (2026); IMDA Digital Economy Report (2025)"
        ]
    },
    {
        "title": "B. Political Factors (P)",
        "subtitle": "",
        "bullets": [
            "OPPORTUNITY (+): National AI Strategy 2.0 (NAIS 2.0) & AI Trailblazers 2.0 initiatives; enterprise AI acceleration frameworks.",
            "THREAT (!): Digital Infrastructure Act compliance layer (2026); 9% GST on import landed unit costs.",
            "MARKETING IMPLICATION: Position G2 as an enterprise-approved hardware platform under Singapore's NAIS 2.0 framework for professional workflows.",
            "Sources: Singapore EDB (2026); Smart Nation SG / MCI NAIS 2.0 (2026)"
        ]
    },
    {
        "title": "B. Economic Factors (E)",
        "subtitle": "",
        "bullets": [
            "OPPORTUNITY (+): Q2 2026 GDP +5.7% YoY led by tech investments; high PMET disposable income.",
            "THREAT (!): Inflation forecast 1.5%–2.5% (MAS); cost-of-living pressures narrowing retail buyers.",
            "MARKETING IMPLICATION: Launch standard SGD price tier alongside a 12-month 0%-interest bank instalment option; introduce B2B enterprise procurement bundles.",
            "Sources: Monetary Authority of Singapore (2026); Ministry of Trade and Industry (2026)"
        ]
    },
    {
        "title": "B. Social Factors (S)",
        "subtitle": "",
        "bullets": [
            "OPPORTUNITY (+): Highly multilingual business pop.; high smartphone penetration.",
            "THREAT (!): AI tech distrust; corporate bans on unauthorized recording devices in law/finance (e.g., Rajah & Tann privacy policies).",
            "TARGET AUDIENCE: Bilingual PMETs (25–45), meeting-heavy roles, privacy-conscious.",
            "MARKETING IMPLICATION: Lead campaigns on LinkedIn with a 'No Camera, Total Confidentiality' message targeting corporate PMETs.",
            "Sources: IMDA Singapore Digital Economy Report (2025); Rajah & Tann Privacy Guidelines (2025)"
        ]
    },
    {
        "title": "B. Technological Factors (T)",
        "subtitle": "",
        "bullets": [
            "OPPORTUNITY (+): Singtel nationwide 5G Standalone (5G SA) coverage (>95%); low-latency AI transmission.",
            "THREAT (!): MemoMind One (US$399–599, camera-free); Meta Ray-Ban Display (US$799, 12MP camera, ~80% market share).",
            "MARKETING IMPLICATION: Own the corporate B2B productivity niche that consumer camera glasses cannot enter due to strict workplace privacy laws.",
            "Sources: Singtel 5G SA Release (2022/2026); Road to VR (2026); Gulf News (2026)"
        ]
    },
    {
        "title": "B. Legal Factors (L)",
        "subtitle": "",
        "bullets": [
            "OPPORTUNITY (+): PDPA provides one clear regulatory framework building consumer/investor trust.",
            "THREAT (!): Extraterritorial consent rules apply to overseas vendors; ad claims must be substantiated by law.",
            "MARKETING IMPLICATION: Appoint a Singapore Data Protection Lead within 90 days of launch; route marketing copy through legal review.",
            "Sources: PDPC Singapore (2025); Hawksford Compliance Guide (2025)"
        ]
    },
    {
        "title": "B. Environmental Factors (E)",
        "subtitle": "",
        "bullets": [
            "OPPORTUNITY (+): Durable titanium/magnesium build extends product life; voluntary trade-in leadership.",
            "THREAT (!): Mandatory NEA Extended Producer Responsibility (EPR) registration for EEE importers.",
            "MARKETING IMPLICATION: Launch 'Even Trade-In' program exceeding mandatory NEA EPR minimums with a public Year-1 take-back target.",
            "Sources: National Environment Agency (2025); Ministry of Sustainability and the Environment (2025)"
        ]
    },
    {
        "title": "B. Synthesis – Priority Matrix",
        "subtitle": "",
        "bullets": [
            "TOP 3 OPPORTUNITIES:",
            "  1. NAIS 2.0 Enterprise Alignment (Political) — Immediate",
            "  2. GDP Growth + Premium PMET Tech Spend (Economic) — Immediate",
            "  3. Multilingual PMET Workplace Fit (Social) — Long-term",
            "TOP 3 THREATS:",
            "  1. PDPA Extraterritorial Data Compliance (Legal) — Immediate",
            "  2. Consumer Smart-Wearables Market Saturation (Technological) — Immediate",
            "  3. Inflation Pressure on Hardware Budgets (Economic) — Long-term"
        ]
    },
    {
        "title": "B. Synthesis – Strategic Recommendations & Expected Outcomes",
        "subtitle": "",
        "bullets": [
            "Enterprise Alignment: Integrate under NAIS 2.0 framework → Accelerates B2B market entry.",
            "Tiered Pricing: Standard RRP + 12-mo 0% bank instalment → Overcomes upfront purchase friction.",
            "Privacy Campaign: LinkedIn messaging focused on 'No Camera' → Converts privacy-conscious PMETs.",
            "Channel Differentiation: Own corporate B2B productivity → Bypasses low-margin consumer price wars.",
            "Compliance & ESG: Local DPO + Branded Trade-in program → Guarantees PDPA/EPR compliance."
        ]
    },
    {
        "title": "References (Harvard Style)",
        "subtitle": "",
        "bullets": [
            "Dealroom.co (2026) Even Realities Company Profile. Available at: app.dealroom.co [Accessed 13 Aug. 2026].",
            "Grand View Research (2026) Smart Glasses Market Report 2026–2030. Available at: grandviewresearch.com [Accessed 13 Aug. 2026].",
            "Infocomm Media Development Authority (2025) Singapore Digital Economy Report 2025. Singapore: IMDA.",
            "Monetary Authority of Singapore (2026) Monetary Policy Statement, July 2026. Available at: mas.gov.sg [Accessed 13 Aug. 2026].",
            "National Environment Agency (2025) Extended Producer Responsibility System for E-waste. Available at: nea.gov.sg [Accessed 13 Aug. 2026].",
            "Personal Data Protection Commission (2025) Advisory Guidelines on Key Concepts in the Personal Data Protection Act. Singapore: PDPC.",
            "Rajah & Tann Singapore (2025) Legal Sector AI Governance and Confidentiality Guidelines. Singapore: Rajah & Tann.",
            "Singtel (2022) Singtel's 5G Standalone Network Surpasses 95% Nationwide Coverage. Available at: singtel.com.",
            "Smart Nation Singapore (2026) National AI Strategy 2.0 (NAIS 2.0) Updates. Available at: smartnation.gov.sg."
        ]
    },
    {
        "title": "Group Declaration",
        "subtitle": "",
        "bullets": [
            "We declare that this submission is our own original work except where standard GenAI tools have been used strictly as outlined in our Member AI Contribution Table.",
            "All information generated or assisted by GenAI has been independently verified against official primary and secondary Singaporean data portals.",
            "Full group accountability is accepted for all contents within this presentation.",
            "",
            "Member Sign-off Signatures:",
            "[ Member 1 Signature ]    [ Member 2 Signature ]    [ Member 3 Signature ]",
            "[ Member 4 Signature ]    [ Member 5 Signature ]"
        ]
    },
    {
        "title": "Group Members & AI Use Table",
        "subtitle": "",
        "bullets": [
            "Member 1 | PESTLE: Background, Political | Tool: ChatGPT | Purpose: SG NAIS 2.0 policies",
            "Member 2 | PESTLE: Economic, Social | Tool: ChatGPT | Purpose: GDP & Inflation analysis",
            "Member 3 | PESTLE: Tech, Legal | Tool: ChatGPT | Purpose: 5G SA scan & PDPA compliance",
            "Member 4 | PESTLE: Environmental | Tool: ChatGPT | Purpose: NEA EPR e-waste scan",
            "Member 5 | PESTLE: Synthesis & Recs | Tool: ChatGPT | Purpose: Priority matrix scaffolding"
        ]
    }
]

# Build Slides
blank_layout = prs.slide_layouts[6]

for data in slides_data:
    slide = prs.slides.add_slide(blank_layout)
    
    # Title Box
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11.7), Inches(1.2))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = data["title"]
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = RGBColor(16, 44, 87)
    
    # Subtitle / Body
    if data["subtitle"]:
        p2 = tf.add_paragraph()
        p2.text = data["subtitle"]
        p2.font.size = Pt(18)
        p2.font.color.rgb = RGBColor(80, 80, 80)
        
    if data["bullets"]:
        body_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
        btf = body_box.text_frame
        btf.word_wrap = True
        
        for idx, bullet in enumerate(data["bullets"]):
            bp = btf.paragraphs[0] if idx == 0 else btf.add_paragraph()
            bp.text = bullet
            bp.font.size = Pt(14 if len(data["bullets"]) > 7 else 16)
            bp.font.color.rgb = RGBColor(40, 40, 40)
            bp.space_after = Pt(6)

prs.save("Even_Realities_G2_Singapore_Entry.pptx")
print("Successfully generated 'Even_Realities_G2_Singapore_Entry.pptx'!")